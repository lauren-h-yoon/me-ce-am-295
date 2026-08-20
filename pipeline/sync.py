#!/usr/bin/env python3
"""
ME/CE/AM 295 — Content sync / normalization.

Walks materials/ (docx, pdf, pptx, md), normalizes each file to clean Markdown
with YAML frontmatter under content/, and maintains content/manifest.json — the
"index" of what exists, which week it belongs to, its type, and access level.

Robustness contract:
  * Source files are hashed (raw bytes). Re-running only re-normalizes files
    whose hash changed; unchanged files are skipped.
  * Files removed/renamed at the source are purged from content/ and the manifest.
  * Instructor-only material (answer keys, rubrics) is NEVER written to content/
    (which feeds the public site + student DB); it is recorded as skipped.

Both the vector index (Slack bot grounding) and the Quarto site read content/,
so they can never drift from each other.

Usage:
    python -m pipeline.sync              # incremental sync
    python -m pipeline.sync --force      # re-normalize everything
    python -m pipeline.sync --dry-run    # show plan, write nothing
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

from pipeline import config as C


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def slugify(text: str) -> str:
    text = re.sub(r"[^A-Za-z0-9]+", "-", text).strip("-").lower()
    return re.sub(r"-{2,}", "-", text) or "untitled"


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def section_for(week: int, doc_type: str) -> str:
    """Where the normalized file lives under content/.

    Assessments (quizzes) go to a separate _assessments/ tree so the weekly
    filesystem agent (scoped to content/week-NN + globals) never sees quiz
    questions. They stay out of the public site too.
    """
    if doc_type == "quiz":
        return f"_assessments/week-{week:02d}" if week and week >= 1 else "_assessments"
    if week and week >= 1:
        return f"week-{week:02d}"
    if doc_type == "syllabus":
        return "syllabus"
    if doc_type == "study_guide":
        return "study-guides"
    if doc_type in ("reference", "reading"):
        return "references"
    return "general"


def should_skip(path: Path) -> str | None:
    """Return a reason string if this path must be skipped, else None."""
    parts = set(path.parts)
    if parts & C.EXCLUDE_DIRS:
        return "excluded_dir"
    name = path.name
    if name.startswith(".") or name.startswith("~$"):
        return "hidden_or_temp"
    if any(m in name for m in C.EXCLUDE_FILE_MARKERS):
        return "superseded_marker"
    if C.DUP_COPY_RE.search(name):
        return "duplicate_copy"
    return None


SUPPORTED = {".md", ".txt", ".docx", ".pdf", ".pptx"}


# ------------------------------------------------------------------
# Format extractors -> Markdown text
# ------------------------------------------------------------------
def extract_docx(src: Path, media_dir: Path) -> str:
    """Convert .docx to GitHub-flavored Markdown via pandoc, extracting media."""
    media_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        ["pandoc", "-f", "docx", "-t", "gfm", "--wrap=none",
         f"--extract-media={media_dir}", str(src)],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        raise RuntimeError(f"pandoc failed: {result.stderr.strip()[:300]}")
    # If pandoc created no media, drop the empty dir to keep the tree clean.
    if media_dir.exists() and not any(media_dir.rglob("*")):
        shutil.rmtree(media_dir, ignore_errors=True)
    return result.stdout


def extract_pdf(src: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(src))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"<!-- page {i} -->\n\n{text}")
    return "\n\n".join(pages)


def extract_pptx(src: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(src))
    out = []
    for n, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {n}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        lines.append(para.text.strip())
        if len(lines) > 1:
            out.append("\n\n".join(lines))
    return "\n\n".join(out)


def extract(src: Path, media_dir: Path) -> str:
    ext = src.suffix.lower()
    if ext in (".md", ".txt"):
        return src.read_text(encoding="utf-8", errors="replace")
    if ext == ".docx":
        return extract_docx(src, media_dir)
    if ext == ".pdf":
        return extract_pdf(src)
    if ext == ".pptx":
        return extract_pptx(src)
    raise ValueError(f"unsupported extension: {ext}")


def title_from(name: str) -> str:
    stem = Path(name).stem
    stem = re.sub(r"^(ME[_ ]?CE[_ ]?(AM[_ ]?)?295[_ ]?)", "", stem, flags=re.I)
    return re.sub(r"[_-]+", " ", stem).strip()


def frontmatter(meta: dict) -> str:
    lines = ["---"]
    for k, v in meta.items():
        if isinstance(v, list):
            lines.append(f"{k}: [{', '.join(v)}]")
        else:
            lines.append(f"{k}: {json.dumps(v) if isinstance(v, str) else v}")
    lines.append("---\n")
    return "\n".join(lines)


# ------------------------------------------------------------------
# Sync
# ------------------------------------------------------------------
def load_manifest() -> dict:
    if C.MANIFEST_PATH.exists():
        return json.loads(C.MANIFEST_PATH.read_text())
    return {"generated_at": None, "files": {}}


def run_sync(force: bool = False, dry_run: bool = False) -> dict:
    manifest = load_manifest()
    old_files: dict = manifest.get("files", {})
    new_files: dict = {}

    added = updated = skipped_unchanged = skipped_excluded = instructor = 0
    seen_sources: set[str] = set()

    sources = sorted(
        p for p in C.MATERIALS_DIR.rglob("*")
        if p.is_file() and p.suffix.lower() in SUPPORTED
    )

    for src in sources:
        rel = str(src.relative_to(C.REPO_ROOT))
        reason = should_skip(src)
        if reason:
            skipped_excluded += 1
            continue
        seen_sources.add(rel)

        name = src.name
        week = C.detect_week(name, str(src.parent))
        doc_type = C.detect_doc_type(name, str(src.parent))
        access = C.access_level_for(doc_type)
        file_hash = sha256_file(src)

        # Instructor-only material never lands in content/.
        if access == C.ACCESS_INSTRUCTOR:
            instructor += 1
            new_files[rel] = {
                "hash": file_hash, "week": week, "doc_type": doc_type,
                "access": access, "status": "skipped_instructor_only",
                "content_path": None,
            }
            continue

        section = section_for(week, doc_type)
        slug = slugify(Path(name).stem)
        content_rel = f"content/{section}/{slug}.md"
        content_path = C.REPO_ROOT / content_rel

        prev = old_files.get(rel)
        if prev and prev.get("hash") == file_hash and content_path.exists() and not force:
            skipped_unchanged += 1
            new_files[rel] = prev
            continue

        if dry_run:
            action = "UPDATE" if prev else "ADD"
            print(f"  [{action}] {rel}  ->  {content_rel}  (week={week}, type={doc_type})")
            (updated if prev else added).__class__  # no-op for type
            if prev:
                updated += 1
            else:
                added += 1
            new_files[rel] = {"hash": file_hash, "week": week, "doc_type": doc_type,
                              "access": access, "status": "planned", "content_path": content_rel}
            continue

        # Normalize.
        media_dir = content_path.parent / "media" / slug
        try:
            body = extract(src, media_dir)
        except Exception as e:
            print(f"  [ERROR] {rel}: {e}", file=sys.stderr)
            continue

        if not body or len(body.strip()) < 40:
            print(f"  [WARN] {rel}: empty/too short after extraction, skipping")
            continue

        meta = {
            "title": title_from(name),
            "week": week,
            "doc_type": doc_type,
            "access": access,
            "agents": C.agents_for(doc_type),
            "source": rel,
            "source_hash": file_hash[:16],
            "normalized_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        }
        # If this file moved to a new content path, remove the stale copy.
        if prev and prev.get("content_path") and prev["content_path"] != content_rel:
            old = C.REPO_ROOT / prev["content_path"]
            if old.exists():
                old.unlink()
                shutil.rmtree(old.parent / "media" / old.stem, ignore_errors=True)

        content_path.parent.mkdir(parents=True, exist_ok=True)
        content_path.write_text(frontmatter(meta) + body.strip() + "\n", encoding="utf-8")

        new_files[rel] = {"hash": file_hash, "week": week, "doc_type": doc_type,
                          "access": access, "status": "normalized", "content_path": content_rel}
        if prev:
            updated += 1
            print(f"  [UPDATE] {rel} -> {content_rel}")
        else:
            added += 1
            print(f"  [ADD]    {rel} -> {content_rel}")

    # Purge content whose source disappeared.
    purged = 0
    for rel, entry in old_files.items():
        if rel not in seen_sources and rel not in new_files:
            cp = entry.get("content_path")
            if cp:
                p = C.REPO_ROOT / cp
                if p.exists() and not dry_run:
                    p.unlink()
                    media = p.parent / "media" / p.stem
                    shutil.rmtree(media, ignore_errors=True)
                print(f"  [PURGE]  {rel} (source removed)")
            purged += 1

    manifest = {
        "generated_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "current_week": C.current_week(),
        "counts": {
            "normalized": sum(1 for f in new_files.values() if f.get("content_path")),
            "instructor_only": instructor,
            "added": added, "updated": updated,
            "skipped_unchanged": skipped_unchanged,
            "skipped_excluded": skipped_excluded,
            "purged": purged,
        },
        "files": new_files,
    }
    if not dry_run:
        C.MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
        C.MANIFEST_PATH.write_text(json.dumps(manifest, indent=2))
        from pipeline.build_memory import build as build_memory
        build_memory()

    return manifest


def main():
    ap = argparse.ArgumentParser(description="Normalize course materials -> content/")
    ap.add_argument("--force", action="store_true", help="re-normalize even if unchanged")
    ap.add_argument("--dry-run", action="store_true", help="show plan, write nothing")
    args = ap.parse_args()

    print(f"Sync: materials={C.MATERIALS_DIR}  ->  content={C.CONTENT_DIR}")
    m = run_sync(force=args.force, dry_run=args.dry_run)
    c = m["counts"]
    print("\nSummary:")
    print(f"  normalized (student-facing) : {c['normalized']}")
    print(f"  added / updated             : {c['added']} / {c['updated']}")
    print(f"  skipped (unchanged)         : {c['skipped_unchanged']}")
    print(f"  skipped (excluded/dup)      : {c['skipped_excluded']}")
    print(f"  instructor-only (not ingested): {c['instructor_only']}")
    print(f"  purged (source removed)     : {c['purged']}")
    print(f"  current teaching week       : {m['current_week']}")


if __name__ == "__main__":
    main()
